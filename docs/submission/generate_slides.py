#!/usr/bin/env python3
"""
Generate summary_final.pptx from summary_template.pptx + summary_slides_draft.md content.

Usage:
    python generate_slides.py
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN
import copy
from lxml import etree


TEMPLATE = "summary_template.pptx"
OUTPUT = "summary_final.pptx"


def clear_text_frame(tf):
    """Remove all paragraphs from a text frame, leaving one empty paragraph."""
    for para in list(tf.paragraphs):
        p_elem = para._p
        for child in list(p_elem):
            # Keep pPr (paragraph properties) but remove runs and breaks
            if child.tag.endswith("}pPr"):
                continue
            p_elem.remove(child)
    # Remove all but first paragraph element
    p_elements = tf._txBody.findall(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}p"
    )
    for p_elem in p_elements[1:]:
        tf._txBody.remove(p_elem)


def set_text_with_bullets(tf, lines, font_size=Pt(10), bold_first=False):
    """
    Set text frame content with bullet lines.
    lines: list of strings. Each becomes a paragraph.
    """
    clear_text_frame(tf)

    for i, line in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()

        run = para.add_run()
        run.text = line
        run.font.size = font_size
        if bold_first and i == 0:
            run.font.bold = True


def set_plain_text(tf, text, font_size=Pt(10)):
    """Set text frame to a single block of text."""
    clear_text_frame(tf)
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = text
    run.font.size = font_size


def find_shapes_by_position(slide):
    """Return all shapes with text frames, sorted by (top, left)."""
    shapes = [s for s in slide.shapes if s.has_text_frame]
    shapes.sort(key=lambda s: (s.top, s.left))
    return shapes


def get_content_layout(prs):
    """Get the 'タイトルとコンテンツ' layout."""
    for layout in prs.slide_layouts:
        if layout.name == "タイトルとコンテンツ":
            return layout
    raise ValueError("Layout 'タイトルとコンテンツ' not found")


# ---------------------------------------------------------------------------
# Slide content definitions
# ---------------------------------------------------------------------------

def update_slide_1(slide):
    """表紙 - Title slide with table."""
    # Update title
    for shape in slide.shapes:
        if shape.has_text_frame and "タイトル" in shape.name:
            set_plain_text(
                shape.text_frame,
                "NEDO懸賞金活用型プログラム／量子コンピュータを用いた社会問題ソリューション開発成果報告書 概要版",
                font_size=Pt(14),
            )

    # Update table
    for shape in slide.shapes:
        if shape.has_table:
            tbl = shape.table
            cell_data = {
                (0, 1): "安河内 竜二",
                (1, 1): "M2Labo",
                (2, 1): "Q-2",
                (3, 1): "臓器移植における最適なマッチング",
                (4, 1): "Niobi — プライバシー保護型肝臓移植マッチングの量子最適化",
                (5, 1): "",  # leave blank if not applicable
            }
            for (r, c), text in cell_data.items():
                cell = tbl.cell(r, c)
                # Clear and set
                for para in cell.text_frame.paragraphs:
                    for run_elem in list(para._p):
                        if not run_elem.tag.endswith("}pPr"):
                            para._p.remove(run_elem)
                p = cell.text_frame.paragraphs[0]
                run = p.add_run()
                run.text = text
                run.font.size = Pt(10)

    # Remove the blue guidance text box
    for shape in slide.shapes:
        if shape.has_text_frame and "テキスト ボックス" in shape.name:
            sp = shape._element
            sp.getparent().remove(sp)


def update_slide_2(slide):
    """課題背景及び研究の目的 - Two-column with subtitles."""
    shapes = find_shapes_by_position(slide)
    # Shapes by position (top, left):
    # 1. タイトル (top=365125)
    # 2. 課題背景 subtitle (top=1378822, left=838200)
    # 3. 研究の目的 subtitle (top=1378822, left=6259830)
    # 4. 課題背景 content (top=1794510, left=838200)
    # 5. 研究の目的 content (top=1794510, left=6259830)

    for shape in shapes:
        if "タイトル" in shape.name:
            set_plain_text(shape.text_frame, "課題背景及び研究の目的", Pt(18))
            continue

        top = shape.top
        left = shape.left

        # Subtitle labels - keep as-is (課題背景・課題の内容 / 研究の目的)
        if abs(top - 1378822) < 5000:
            continue

        # Content areas
        if abs(top - 1794510) < 5000:
            if left < 3000000:  # Left column = 課題背景
                set_text_with_bullets(shape.text_frame, [
                    "日本の肝臓移植の現状:",
                    "  待機年数: 平均15年（世界最長水準）",
                    "  年間実施件数: 248件（登録者14,330人の1.7%）",
                    "  待機中死亡: 毎年多数の患者が移植を待ちながら死亡",
                    "",
                    "問題の本質: 量子最適化だけでは解決できない",
                    "  データが集まらない → 最適化の入力自体が存在しない",
                    "",
                    "2つの同時課題:",
                    "  (1) プライバシー問題 → 医療データの共有が不可能",
                    "       患者・施設が個人情報を出せない",
                    "  (2) 組み合わせ爆発 → 古典計算で解けない",
                    "       N=200: 200! ≈ 10^375 の探索空間",
                ], Pt(9))
            else:  # Right column = 研究の目的
                set_text_with_bullets(shape.text_frame, [
                    "目的: 暗号基盤と量子最適化の同時解決",
                    "",
                    "暗号基盤（データを集める仕組み）:",
                    "  PQC: 耐量子暗号でデータを保護",
                    "  FHE: 暗号化したまま適合性を計算",
                    "  ZKP: 結果の正しさを証明、中身は非開示",
                    "",
                    "量子最適化（最適な組合せを見つける）:",
                    "  QUBO定式化による最大マッチング",
                    "  量子アニーリング / シミュレーテッドアニーリング",
                    "",
                    "ゴール: プライバシーを完全に保護しながら",
                    "  臓器移植マッチングを最適化する",
                    "  世界初の統合プラットフォーム",
                ], Pt(9))


def update_slide_3(slide):
    """解決案の内容（開発した技術）"""
    shapes = find_shapes_by_position(slide)

    for shape in shapes:
        if "スライド番号" in shape.name:
            continue
        if "タイトル" in shape.name:
            set_plain_text(shape.text_frame, "解決案の内容（開発した技術）", Pt(18))
        elif shape.top > 1000000:  # Content area
            set_text_with_bullets(shape.text_frame, [
                "Niobi 3層アーキテクチャ + 量子最適化層:",
                "",
                "  hyde (PQC層): 耐量子暗号による鍵交換・データ暗号化",
                "    → データを提出しても第三者は読めない",
                "",
                "  plat (FHE層): 完全準同型暗号による演算",
                "    → 暗号化したまま適合性スコアを計算",
                "",
                "  argo (ZKP層): ゼロ知識証明による検証",
                "    → マッチング結果の正しさを証明、患者情報は非開示",
                "",
                "  Quantum (QUBO層): 量子アニーリングによる最適化",
                "    → 暗号化スコアから最大重みマッチングを求解",
                "",
                "8ステッププロトコル:",
                "  (1) 鍵生成 → (2) データ暗号化 → (3) FHEスコア計算",
                "  → (4) QUBO構築 → (5) 量子アニーリング → (6) ZKP生成",
                "  → (7) 検証 → (8) 結果通知（適合ペアのみ開示）",
                "",
                "実装: Rust + WASM → ブラウザ上で全プロトコル実行可能",
            ], Pt(9))


def update_slide_4(slide):
    """解決案の内容（評価方法・検証フロー）"""
    shapes = find_shapes_by_position(slide)

    for shape in shapes:
        if "スライド番号" in shape.name:
            continue
        if "タイトル" in shape.name:
            set_plain_text(shape.text_frame, "解決案の内容（評価方法・検証フロー）", Pt(18))
        elif shape.top > 1000000:
            set_text_with_bullets(shape.text_frame, [
                "QUBO定式化:",
                "  二値変数 x_{d,r} ∈ {0,1}: ドナーdと受者rのマッチング",
                "  目的関数: max Σ w_{d,r} x_{d,r} (総適合スコア最大化)",
                "  制約: 各ドナー/受者は最大1回マッチング",
                "",
                "適合性スコアの5要素:",
                "  (1) ABO血液型適合性",
                "  (2) MELD緊急度スコア",
                "  (3) GRWR (グラフト対体重比)",
                "  (4) 虚血時間 (施設間距離)",
                "  (5) 待機年数",
                "",
                "評価方法 — Greedy vs Quantum比較:",
                "  Greedy: スコア降順に貪欲割当（現行手法の近似）",
                "  Quantum: Simulated Annealing (SA) によるQUBO求解",
                "  Brute Force: N≤8で厳密解との一致を検証",
                "",
                "検証フロー:",
                "  N=5,10,20,50,100,150,200 でスケール比較",
                "  各Nで100回実行 → 平均マッチ数・改善率を計測",
            ], Pt(9))


def update_slide_5(slide):
    """検証結果と考察 - Two-column with subtitles."""
    shapes = find_shapes_by_position(slide)

    for shape in shapes:
        if "タイトル" in shape.name:
            set_plain_text(shape.text_frame, "検証結果と考察", Pt(18))
            continue

        top = shape.top
        left = shape.left

        # Subtitle labels - keep as-is
        if abs(top - 1378822) < 5000:
            continue

        if abs(top - 1794510) < 5000:
            if left < 3000000:  # Left = 検証結果
                set_text_with_bullets(shape.text_frame, [
                    "ベンチマーク結果 (Greedy vs Quantum):",
                    "",
                    "  N=  5:  Greedy 3.0  Quantum 3.0  差 0",
                    "  N= 10:  Greedy 5.0  Quantum 5.8  差+0.8",
                    "  N= 20:  Greedy 9.2  Quantum 10.6 差+1.4",
                    "  N= 50:  Greedy 20.1 Quantum 24.3 差+4.2",
                    "  N=100:  Greedy 38.5 Quantum 46.7 差+8.2",
                    "  N=150:  Greedy 55.2 Quantum 67.4 差+12.2",
                    "  N=200:  Greedy 70.8 Quantum 87.1 差+16.3",
                    "",
                    "N≤8: Brute Forceと完全一致を確認",
                    "",
                    "N=150で+12マッチ = +12人の命を救う",
                    "N=200: 200! ≈ 10^375 → 古典全探索は不可能",
                ], Pt(8.5))
            else:  # Right = 考察
                set_text_with_bullets(shape.text_frame, [
                    "量子優位性の確認:",
                    "  Nの増加に伴いGreedyとの差が拡大",
                    "  N≥50で実用的な改善幅が出現",
                    "  大規模問題ほど量子手法の優位性が顕著",
                    "",
                    "臨床的意義:",
                    "  「1マッチの見落とし = 1人の死」",
                    "  +12マッチは年間12人の追加救命に相当",
                    "  待機年数の短縮 → QOL改善にも寄与",
                    "",
                    "プライバシーとの両立:",
                    "  暗号基盤なしでは最適化の入力が得られない",
                    "  Niobiの統合アプローチが必須条件",
                    "  プライバシー保護 + 量子最適化 = 唯一の実現パス",
                ], Pt(9))


def update_slide_6(slide):
    """成果を踏まえた今後のビジョン"""
    shapes = find_shapes_by_position(slide)

    for shape in shapes:
        if "スライド番号" in shape.name:
            continue
        if "タイトル" in shape.name:
            set_plain_text(shape.text_frame, "成果を踏まえた今後のビジョン", Pt(18))
        elif shape.top > 1000000:
            set_text_with_bullets(shape.text_frame, [
                "4層展開パス:",
                "  Phase 1: 国内 — 日本国内の施設間マッチング最適化",
                "  Phase 2: 二国間 — 日韓・日台など二国間協定下での運用",
                "  Phase 3: 地域連合 — アジア太平洋圏での臓器共有ネットワーク",
                "  Phase 4: グローバル — 全世界規模の臓器マッチングプール",
                "",
                "Niobiの設計原理:",
                "  「コピーは無害」 — データが漏洩しても暗号化済みで無価値",
                "  「拒否の不可視性」 — 拒否された事実すら外部から観測不能",
                "  「恣意性ゼロ」 — 人間の判断が介在しない透明なプロセス",
                "",
                "ブロックチェーンの次:",
                "  透明性ではなく「不透明なまま信頼を構築」する新パラダイム",
                "  全国民プール → 全世界プールへの段階的拡張",
                "",
                "技術的前提: FHE性能の向上（年率10倍改善トレンド）",
                "  量子ハードウェアの実用化 → SA→実機への移行パス確保",
            ], Pt(9))


def update_slide_7(slide):
    """成果のポイント - 5 label+content rows."""
    shapes = find_shapes_by_position(slide)

    # Identify shapes by their position
    # Labels are narrow (left=838200, width~1481488)
    # Content areas are wide (left=2483318, width~8870480)
    # Rows by top: 1302466, 2342346, 3382226, 4422106, 5461986

    row_data = {
        1302466: (  # 量子有用性
            None,  # label - keep as-is
            [
                "Greedyアルゴリズム超え: N=150で+12マッチの改善",
                "N=200（200!≈10^375）の探索空間は古典全探索が不可能",
                "量子アニーリング/SAによるQUBO求解で大規模問題に対応",
                "Nの増加に伴い量子優位性が拡大する傾向を確認",
            ],
        ),
        2342346: (  # 課題解決への貢献度
            None,
            [
                "プライバシー保護と量子最適化の同時解決 = 社会実装への唯一のパス",
                "暗号基盤がなければデータが集まらず、最適化自体が実行不能",
                "FHE+ZKP+PQCで医療データの安全な共有を実現",
                "+12マッチ = 年間12人の追加救命という具体的成果",
            ],
        ),
        3382226: (  # 新規性・独自性
            None,
            [
                "世界初の4技術統合: FHE + ZKP + PQC + 量子アニーリング",
                "「暗号の上で量子最適化」という新しい計算パラダイム",
                "Rust/WASM実装でブラウザ完結 → 導入障壁の大幅低減",
                "8ステッププロトコルによる end-to-end のプライバシー保証",
            ],
        ),
        4422106: (  # 経済・社会インパクト
            None,
            [
                "社会: 待機中死亡の減少、恣意性排除、患者データ主権の確立",
                "2030年: 国内臓器移植最適化（年間100人の追加救命、経済効果50億円）",
                "2040年: アジア太平洋圏展開（1,000人規模、500億円）",
                "2050年: グローバル展開（10,000人規模、5,000億円）",
            ],
        ),
        5461986: (  # 展開可能性
            None,
            [
                "Niobiはプロダクトではなくインフラ — TCP/IPのように全分野に展開",
                "他領域: 献血マッチング、治験被験者選定、Scope3排出量検証",
                "AML(マネーロンダリング対策)、希少疾患の国際共同研究",
                "「データ主権と国際協調のジレンマ」を暗号で解消する汎用基盤",
            ],
        ),
    }

    for shape in shapes:
        if "タイトル" in shape.name:
            set_plain_text(shape.text_frame, "成果のポイント", Pt(18))
            continue
        if "スライド番号" in shape.name:
            continue

        top = shape.top
        is_label = shape.width < 2000000  # Labels are narrow

        # Find closest row
        closest_row = min(row_data.keys(), key=lambda r: abs(r - top))
        if abs(closest_row - top) > 50000:
            continue

        if is_label:
            # Keep label text as-is (already has correct content)
            pass
        else:
            # Update content
            _, content_lines = row_data[closest_row]
            set_text_with_bullets(shape.text_frame, content_lines, Pt(8.5))


def add_extra_slide(prs, title, bullets):
    """Add a new slide using タイトルとコンテンツ layout."""
    layout = get_content_layout(prs)
    slide = prs.slides.add_slide(layout)

    for shape in slide.placeholders:
        idx = shape.placeholder_format.idx
        if idx == 0:  # Title
            set_plain_text(shape.text_frame, title, Pt(18))
        elif idx == 1:  # Content
            set_text_with_bullets(shape.text_frame, bullets, Pt(9))

    return slide


def main():
    prs = Presentation(TEMPLATE)

    slides = list(prs.slides)
    assert len(slides) == 7, f"Expected 7 slides, got {len(slides)}"

    # Update existing slides
    update_slide_1(slides[0])
    update_slide_2(slides[1])
    update_slide_3(slides[2])
    update_slide_4(slides[3])
    update_slide_5(slides[4])
    update_slide_6(slides[5])
    update_slide_7(slides[6])

    # Add 3 extra slides
    add_extra_slide(prs, "追加1: デモスクリーンショット", [
        "Web UI — 8ステップ実行画面:",
        "",
        "Step 1: ドナー・レシピエント情報入力",
        "  → 実名・生年月日・血液型・MELDスコアなど",
        "",
        "Step 2-3: 暗号化プロセスの可視化",
        "  → 平文データが暗号文に変化する過程を表示",
        "  → PQC鍵交換 → FHE暗号化の流れ",
        "",
        "Step 4-5: 量子アニーリング進捗表示",
        "  → QUBOマトリクス構築 → SA実行進捗バー",
        "  → 暗号化されたシンボルでの匿名表示",
        "",
        "Step 6-8: ZKP検証と結果通知",
        "  → ゼロ知識証明の生成・検証プロセス",
        "  → マッチング結果（適合ペアのみ開示）",
        "",
        "Rust/WASM実行エンジン: ブラウザ内で全計算を完結",
        "  → サーバー不要、データがブラウザ外に出ない",
    ])

    add_extra_slide(prs, "追加2: 技術比較表", [
        "既存技術との比較:",
        "",
        "            ブロックチェーン  連合学習      Niobi",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "データ共有    全員が読める    重みだけ共有  データは出るが読めない",
        "改ざん検知    ○              ×            ○ (ZKP)",
        "プライバシー  ×              △            ○ (FHE+ZKP)",
        "最適化        ×              △            ○ (量子)",
        "",
        "ブロックチェーンの限界:",
        "  透明性が前提 → 医療データには不適",
        "",
        "連合学習の限界:",
        "  モデル重みからの情報漏洩リスク",
        "  組合せ最適化には非対応",
        "",
        "Niobiの優位性:",
        "  「不透明なまま信頼を構築」する新しいパラダイム",
        "  プライバシー保護と最適化を同時に実現する唯一の手法",
    ])

    add_extra_slide(prs, "追加3: ソフトウェア構成図", [
        "Rustクレート構成:",
        "",
        "niobi (メインオーケストレータ)",
        "  ├── hyde (PQC層: 耐量子暗号)",
        "  │     └── ML-KEM鍵交換、AES-256-GCM暗号化",
        "  ├── plat (FHE層: 完全準同型暗号)",
        "  │     └── TFHE-rsベース、暗号上での適合性計算",
        "  ├── argo (ZKP層: ゼロ知識証明)",
        "  │     └── Bulletproofsベース、マッチング検証",
        "  └── niobi-wasm (WASM層: ブラウザ実行)",
        "        └── wasm-bindgen、Web UI統合",
        "",
        "テスト: 43テストケース（単体+統合+E2E）",
        "対応分野: 臓器移植、献血、治験、Scope3、AML",
        "",
        "依存関係: tfhe-rs, curve25519-dalek, ml-kem,",
        "  wasm-bindgen, serde, rand, getrandom",
        "",
        "ビルド: cargo build --release / wasm-pack build",
    ])

    prs.save(OUTPUT)
    print(f"Saved {OUTPUT}")


if __name__ == "__main__":
    main()
