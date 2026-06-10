#!/usr/bin/env python3
"""
Generate summary_official.pptx from the OFFICIAL summary template
(research_result_report_summary_templete.pptx, copied locally) + report content.

- Slide 1: fill 表紙 table (PROVISIONAL name/team pending 事務局 confirmation;
  set NAME/TEAM to "" to anonymize), delete the 作成要領 instruction box.
- Slides 2-7: replace each blue guidance placeholder with concise bullets
  sourced from report_draft.md. Labels are kept. Body font >= 10pt, blind.
Content is identified by shape_id (stable across the template).
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BASE = Path(__file__).parent
TEMPLATE = BASE / "summary_template_official.pptx"
OUTPUT = BASE / "summary_official.pptx"

JP = "ＭＳ Ｐゴシック"
BLACK = RGBColor(0x1A, 0x1A, 0x1A)

# 概要版は完全匿名（事務局確認済: 氏名欄は削除）。氏名・団体名は空欄にする。
COVER = {
    "応募代表者氏名": "",
    "チーム名": "",
    "選択課題ID": "Q-2",
    "選択課題名": "創薬エコシステムの強化に向けた医療データ共有アプリケーション・アルゴリズムの開発",
    "提案名": "Niobi — プライバシー保護型臓器移植マッチングの分散最適化",
    "解決案ID": "",
}

# (text, indent_level). Use ** for bold. Lines starting with no bullet if level<0.
S2_LEFT = [
    ("医療データは『出せない』ために最適化できない。創薬では薬剤治療データ（治験・失敗例）が企業ごとに分断、臓器移植では病院が患者データを共有しない — **共有できないことが根本原因**。本基盤は汎用で、最も切迫した臓器移植を最初の実証、本命適用先を創薬データ共有に置く。", 0),
    ("**データ共有を阻む3層の壁:**", 0),
    ("国内 — 高粒度データ(HLA完全型等)が各施設に滞留。追加収集に再同意が必要で制度的摩擦", 1),
    ("国際 — GDPR/HIPAA/APPIにより越境医療データ共有が法的に不可能", 1),
    ("潜在ドナー — 意思表示の安全な登録・管理基盤が不在", 1),
    ("**割当(単一/多臓器)は古典で最適に解ける**が、直接マッチング不成立の不適合ペアの**交換(2-/3-way cycle-cover)**は頂点素な最大重み集合パッキング=NP困難。貪欲は構造的に劣り大域最適化が要る", 0),
]
S2_RIGHT = [
    ("「個人情報を保護する」だけでなく**「共有時点で受領者にとっての識別性を暗号学的に消す」**（受領者視点。元の管理者には個人情報として適用継続）", 0),
    ("暗号を**計算層と保護層に分離**し、計算層を完全に破っても個人情報に到達不能な設計原則を確立", 0),
    ("hyde暗号文は鍵なしでは乱数と区別不能 → 再識別手段を持たない受領者の視点では「個人データ」に該当しないと合理的に主張可（ECJ SRB判決の射程、元管理者の義務は残る）", 0),
    ("**データ共有の制度的障壁そのものを解消**し、その上で量子最適化を適用する", 0),
]
S3 = [
    ("**Niobi — データフローに沿った5層アーキテクチャ**", 0),
    ("(0) 認証ゲート層 — TEE(Intel TDX/AMD SEV-SNP)内にFHE鍵を封入。TPM署名＋病院署名で認証済データのみFHE暗号化。**病院はFHEライブラリ不要**", 0),
    ("(1) 計算層 — plat(Threshold FHE + Bootstrapping)で暗号状態のまま適合率計算。秘密鍵はt-of-n閾値分割(独立機関のTEE内)", 0),
    ("(2) 保護層 — hyde(PQC/ML-KEM-768)で個人情報(氏名・連絡先)を暗号学的に分離。鍵はTPMに紐づき個人が保管", 0),
    ("(3) 検証層 — argo(ZKP)で計算の正当性を中身非開示のまま証明", 0),
    ("(4) 最適化層 — QUBO定式化で交換cycle-cover(NP困難)を最適化(多臓器割当は適合スコアの応用層)", 0),
    ("**全構成要素は今日の技術で実現**(NIST標準PQC／FHE実装＋GPU加速／古典QUBO)。将来の量子を前提としない。production_8192で50.91ms/ペア実測、MKFHE bootstrappingで異鍵暗号文の比較を復号なしで実証", 0),
]
S4 = [
    ("**評価指標:** マッチ数・適合スコア・計算時間を既存手法と比較", 0),
    ("**ベースライン(単一臓器二部マッチング):** Greedy法 / Hungarian法(O(N³)最適解) / Brute Force(N≤8で正解一致を検証)", 0),
    ("**交換cycle-cover:** ランダム不適合プールで貪欲サイクル検出 vs QUBO(焼きなまし) vs 厳密最適を移植数で比較(各サイズ8シード、2-/3-way)。多臓器割当はOPTN/SRTR 2023文献値で適合スコアを検証", 0),
    ("**暗号コスト:** plat crateでkeygen/encrypt/score/decryptを実測(test_small / research_2048 / production_8192)", 0),
    ("**GPU:** CUDA/wgpu NTTバックエンドをRTX 3090で実測", 0),
    ("**実装:** Rust(--release)。D-Wave SimulatedAnnealingSamplerでも同等結果を再現確認", 0),
]
S5_LEFT = [
    ("**単一臓器:** Hungarianが全スケールで最適、QUBO(焼きなまし)は2-6%劣る(正直な評価)。N=200でGreedy比+12マッチ", 0),
    ("**多臓器(負の結果):** 救命数では分離可能=smart greedy=QUBO=厳密最適、量子優位なし。**交換cycle-cover(本命):** QUBOが貪欲に**8/8勝利**・厳密最適に一致・差は規模で拡大(+1.5→+3.8移植)", 0),
    ("**FHE実測:** composite_score 21µs、full pipeline **50.91ms/ペア**(production_8192)", 0),
    ("**スケール:** 全臓器スコアリング国内11分→**GPU 48秒**、全世界18分(冷阻血時間内)", 0),
    ("**GPU:** N=8192でCUDA **14x**高速化(暗号安全性は不変)", 0),
]
S5_RIGHT = [
    ("量子の有用性は計算速度ではなく、**社会実装の前提条件(プライバシー保護)と組み合わせて初めて意味を持つ**", 0),
    ("量子最適化だけでは解けない — 病院がデータを出さないから。hyde/argo/platで初めて量子に渡すデータが揃う", 0),
    ("割当(単一/多臓器)は古典で十分。**真にNP困難な交換cycle-coverでQUBOが必要** — 貪欲は容易な2-wayを取り優れた3-wayを潰す", 0),
    ("MKFHE(理論最強)はKEPスケールで非両立と判明 → **Threshold FHE + hyde 2層分離**へ設計転換", 0),
]
S6 = [
    ("**「FTQCを待たずに社会実装を開始し、FTQCの完成とともにスケールする」**設計", 0),
    ("現在のNISQ/アニーリングでNiobiの基本構造は動作。全構成要素は今日存在する技術", 0),
    ("FTQC(2030-2040)完成で: 量子実機のQUBO求解がシミュレータを凌駕／N≥10,000の厳密解が到達可能", 0),
    ("**4層展開:** 国内プール(JOTN統合) → 二国間(日米独, Threshold FHE閾値を国際機関で構成) → 地域連合(アジア・EU) → グローバルプール", 0),
    ("遺伝学展開(80億規模)では最適化が組み合わせ的に爆発し古典厳密解が破綻、**量子アニーリングがスケーラブルな求解経路**になる(forward-looking)", 0),
    ("**ロードマップ:** 2026原理検証／2028パイロット／2029国内／2030二国間／2040グローバル／2050全分野", 0),
]
# Slide 7: 5 short boxes (keep to ~2 lines each)
S7_QUANTUM = [
    ("**今までになかった量子有用性を創出する問題構造の発見**：割当は分離可能で量子不要(負の結果)→**交換cycle-cover(NP困難・APX困難)はQUBOが効く**(正の結果)。貪欲に8/8勝・厳密最適一致・差は規模で拡大。問題のカテゴリを正しく選び直した", 0),
]
S7_NOVELTY = [
    ("**計算層と保護層の分離原則**(計算層を破っても個人情報に到達不能)／MKFHE限界の特定→Threshold FHE+hyde転換／FHE・MKFHE bootstrapping Rust実装(GPU14x)／TEE認証ゲート型暗号化／Threshold FHE+ZKP+PQC+TEE+交換cycle-cover QUBO統合", 0),
]
S7_CONTRIB = [
    ("Q-2の根本障壁=医療データ共有のプライバシー問題を計算層・保護層の分離で解消し量子最適化を適用。**身元情報の鍵は個人がTPM保持／計算の鍵はt-of-n閾値機関に分散＝単一機関に鍵が集中しない**。創薬(治験データ企業横断共有・患者リクルート最適化・規制連携)へ同一基盤で展開", 0),
]
S7_IMPACT = [
    ("**社会:** 待機中死亡の減少、手続き的恣意性の排除と評価関数の透明化、データ主権を個人へ。基盤はQoL領域横断（創薬効率化・治験・希少疾患・遺伝学）に展開。 **経済(臓器移植の実証規模):** 2030 +44.5〜67億円/年、2040 +2,000億円/年、2050 +7,500億円/年(全臓器マッチ効率+10〜15%、移植1件≈5,000万円)", 0),
]
S7_EXPAND = [
    ("hyde/plat/argoは汎用インフラ。献血・国際治験・Scope3 CO₂・AML・希少疾患・創薬など「データを出せないから最適化できない」構造の問題すべてに適用。データ主権は常に個人に", 0),
]

# slide_index (0-based) -> { shape_id: (lines, font_pt) }
FILL = {
    1: {3: (S2_LEFT, 11), 4: (S2_RIGHT, 11)},
    # slide index 2 (解決案の内容/開発した技術) is built as a diagram, not bullets
    3: {3: (S4, 11)},
    4: {3: (S5_LEFT, 11), 4: (S5_RIGHT, 11)},
    5: {3: (S6, 11.5)},
    6: {7: (S7_QUANTUM, 10.5), 3: (S7_CONTRIB, 10.5), 12: (S7_NOVELTY, 10.5),
        11: (S7_IMPACT, 10.5), 5: (S7_EXPAND, 10.5)},
}


def _apply_run(run, size, bold):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = JP
    run.font.color.rgb = BLACK
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {}); rPr.append(ea)
    ea.set("typeface", JP)


def _clear_tf(tf):
    txbody = tf._txBody
    for p in txbody.findall(qn("a:p")):
        txbody.remove(p)


def _add_line(tf, text, level, size):
    p = tf._txBody.makeelement(qn("a:p"), {})
    tf._txBody.append(p)
    pPr = p.makeelement(qn("a:pPr"), {})
    pPr.set("lvl", str(level))
    # bullet for level 0/1, none otherwise handled by template; keep template bullets
    p.append(pPr)
    from pptx.text.text import _Paragraph
    para = _Paragraph(p, tf)
    # parse **bold**
    import re
    for tok in re.split(r"(\*\*[^*]+\*\*)", text):
        if not tok:
            continue
        bold = tok.startswith("**") and tok.endswith("**")
        run = para.add_run()
        run.text = tok[2:-2] if bold else tok
        _apply_run(run, size, bold)


def fill_text_frame(tf, lines, size):
    _clear_tf(tf)
    for text, level in lines:
        _add_line(tf, text, level, size)
    tf.word_wrap = True


# ---------------------------------------------------------------------------
# Slide 3 architecture diagram
# ---------------------------------------------------------------------------

def _box(slide, x, y, w, h, title, body, fill, tcol=RGBColor(0xFF, 0xFF, 0xFF),
         tsz=11, bsz=10):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF); sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = Emu(20000); tf.margin_bottom = Emu(20000)
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    r.font.size = Pt(tsz); r.font.bold = True; r.font.name = JP; r.font.color.rgb = tcol
    if body:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run(); r2.text = body
        r2.font.size = Pt(bsz); r2.font.name = JP; r2.font.color.rgb = tcol
    return sp


def _arrow(slide, cx, top, h, label=None):
    a = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                               Inches(cx - 0.13), Inches(top), Inches(0.26), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = RGBColor(0x9A, 0x9A, 0x9A)
    a.line.fill.background(); a.shadow.inherit = False
    if label:
        tb = slide.shapes.add_textbox(Inches(cx + 0.25), Inches(top - 0.04),
                                      Inches(4.6), Inches(h + 0.08))
        tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = label
        r.font.size = Pt(10); r.font.name = JP; r.font.italic = True
        r.font.color.rgb = RGBColor(0x55, 0x55, 0x55)


def build_slide3_diagram(slide):
    # empty the content placeholder so only the diagram shows
    for sh in slide.shapes:
        if sh.shape_id == 3 and sh.has_text_frame:
            _clear_tf(sh.text_frame)
            sh.text_frame._txBody.append(sh.text_frame._txBody.makeelement(qn("a:p"), {}))

    BLUE0 = RGBColor(0x2E, 0x5B, 0x88)
    BLUE1 = RGBColor(0x37, 0x7D, 0xA6)
    TEAL = RGBColor(0x4A, 0x90, 0xA4)
    GREEN = RGBColor(0x4E, 0x8A, 0x5B)
    GRAY = RGBColor(0x60, 0x60, 0x60)
    ORANGE = RGBColor(0xC0, 0x62, 0x2C)
    LIGHT = RGBColor(0xF0, 0xF0, 0xF0)

    lx, lw = 0.92, 7.0
    cx = lx + lw / 2
    _box(slide, lx, 1.42, lw, 0.5, "医療データ（病院・患者）",
         "平文・各施設に滞留。GDPR/HIPAA/APPIで越境共有が不可能", GRAY, bsz=10)
    _arrow(slide, cx, 1.96, 0.18, "TPM署名＋病院署名で認証")
    _box(slide, lx, 2.16, lw, 0.6, "(0) 認証ゲート層 — TEE (Intel TDX / AMD SEV-SNP)",
         "FHE公開鍵をTEE内に封入し、認証済データのみFHE暗号化。病院はFHEライブラリ不要", BLUE0)
    _arrow(slide, cx, 2.80, 0.18, "FHE暗号文")
    _box(slide, lx, 3.00, lw, 0.62, "(1) 計算層 — plat (Threshold FHE + Bootstrapping)",
         "暗号状態のまま適合率計算／秘密鍵は t-of-n 閾値分割（独立機関のTEE内）", BLUE1)
    _arrow(slide, cx, 3.66, 0.18, "匿名スコアグラフ（個人情報を含まない）")
    _box(slide, lx, 3.86, lw, 0.55, "(3) 検証層 — argo (ZKP)",
         "計算の正当性を中身非開示のまま証明", TEAL)
    _arrow(slide, cx, 4.45, 0.16)
    _box(slide, lx, 4.63, lw, 0.6, "(4) 最適化層 — QUBO",
         "交換cycle-cover最適化（NP困難・集合パッキング）。多臓器割当は適合スコア応用層", GREEN)
    _arrow(slide, cx, 5.27, 0.16)
    _box(slide, lx, 5.45, lw, 0.45, "最適マッチ結果（1ドナーから最大8人を救命）", "", GRAY, tsz=11)

    # right column: 保護層 hyde (parallel) + tech note
    rx, rw = 8.35, 4.05
    _box(slide, rx, 2.16, rw, 1.46,
         "(2) 保護層 — hyde (PQC / ML-KEM-768)",
         "個人情報（氏名・連絡先）を計算層から暗号学的に分離。鍵はTPMに紐づき個人が保管。"
         "暗号文は鍵なしでは乱数と区別不能＝受領者視点で識別性なし（元管理者には適用継続）", ORANGE, bsz=10)
    _box(slide, rx, 3.86, rw, 1.37,
         "全構成要素は今日の技術で実現",
         "NIST標準PQC／FHE実装＋GPU加速(14x)／古典QUBO。将来の量子を前提としない。"
         "production_8192: 50.91ms/ペア実測。MKFHE bootstrappingで異鍵暗号文の比較を復号なしで実証",
         LIGHT, tcol=RGBColor(0x22, 0x22, 0x22), bsz=10)


# ---------------------------------------------------------------------------
# Free slides (8-10)
# ---------------------------------------------------------------------------

def _new_slide(prs, layout, title):
    slide = prs.slides.add_slide(layout)
    # set title, remove other placeholders so no prompt text shows
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx == 0 or ph == slide.shapes.title:
            continue
        ph._element.getparent().remove(ph._element)
    t = slide.shapes.title
    t.text = title
    t.text_frame.word_wrap = True
    for p in t.text_frame.paragraphs:
        for r in p.runs:
            r.font.name = JP; r.font.bold = True; r.font.size = Pt(30)
    return slide


def _cell(cell, text, sz=10.5, bold=False, fill=None, color=RGBColor(0x22, 0x22, 0x22)):
    if fill is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_top = Emu(10000); cell.margin_bottom = Emu(10000)
    cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.08)
    tf = cell.text_frame; tf.word_wrap = True
    _clear_tf(tf)
    p = tf._txBody.makeelement(qn("a:p"), {}); tf._txBody.append(p)
    from pptx.text.text import _Paragraph
    para = _Paragraph(p, tf)
    r = para.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold; r.font.name = JP; r.font.color.rgb = color


def build_comparison_slide(prs, layout):
    slide = _new_slide(prs, layout, "既存手法との比較 — なぜNiobiか")
    cols = ["評価軸", "ブロックチェーン", "連合学習", "TEE単独", "Niobi"]
    rows = [
        ("データ共有の仕組み", "全員が平文を閲覧", "勾配・重みのみ共有", "TEE内で平文に復元", "データは出るが鍵なしでは読めない"),
        ("プライバシー保護", "✕ 公開台帳", "△ 勾配から復元リスク", "△ ハード信頼・サイドチャネル", "◎ FHE＋ZKP＋PQC"),
        ("改ざん検知・正当性", "◎ 台帳で保証", "✕", "△ リモート認証", "◎ argo(ZKP)＋台帳監査"),
        ("組合せ最適化", "✕", "△ 限定的", "○", "◎ QUBO(交換cycle-cover)"),
        ("鍵管理主権", "各自が保持", "サーバに集約", "ハード製造者依存", "個人(TPM)＋t-of-n閾値"),
    ]
    nrow, ncol = len(rows) + 1, len(cols)
    gtbl = slide.shapes.add_table(nrow, ncol, Inches(0.55), Inches(1.55),
                                  Inches(12.2), Inches(4.6))
    tbl = gtbl.table
    tbl.columns[0].width = Inches(2.3)
    for c in range(1, 4):
        tbl.columns[c].width = Inches(2.55)
    tbl.columns[4].width = Inches(2.25)
    HEAD = RGBColor(0x22, 0x33, 0x44)
    NIOBI = RGBColor(0xE7, 0xF0, 0xD9)
    for c, name in enumerate(cols):
        fill = RGBColor(0x3A, 0x6B, 0x35) if name == "Niobi" else HEAD
        _cell(tbl.cell(0, c), name, sz=11, bold=True, fill=fill, color=RGBColor(0xFF, 0xFF, 0xFF))
    for ri, row in enumerate(rows, 1):
        _cell(tbl.cell(ri, 0), row[0], sz=10.5, bold=True, fill=RGBColor(0xF0, 0xF0, 0xF0))
        for ci in range(1, 5):
            fill = NIOBI if ci == 4 else None
            bold = ci == 4
            _cell(tbl.cell(ri, ci), row[ci], sz=10, bold=bold, fill=fill)
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(6.3), Inches(12.2), Inches(0.6))
    tf = tb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = "いずれも「データを出せないから最適化できない」構造を解けない。Niobiは唯一、4軸（共有・保護・正当性・最適化）を同時に満たす。"
    r.font.size = Pt(11); r.font.bold = True; r.font.name = JP; r.font.color.rgb = RGBColor(0x3A, 0x6B, 0x35)


def build_benchmark_slide(prs, layout):
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    slide = _new_slide(prs, layout, "検証結果 — ベンチマークと量子有用性")
    cd = CategoryChartData()
    cd.categories = ["8", "12", "16", "20", "30"]
    cd.add_series("貪欲 cycle検出", (5.2, 7.8, 11.2, 13.8, 22.5))
    cd.add_series("QUBO(焼きなまし)", (6.8, 10.0, 13.9, 17.4, 26.2))
    gframe = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                                    Inches(0.5), Inches(1.5), Inches(6.7), Inches(4.0), cd)
    chart = gframe.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.has_title = True
    chart.chart_title.text_frame.text = "交換cycle-cover: 移植数（不適合ペア数N別, 8シード平均）"
    for s in chart.plots[0].series:
        pass
    try:
        for ax in (chart.category_axis, chart.value_axis):
            for r in ax.tick_labels.font._element.getparent().iter():
                pass
    except Exception:
        pass
    cap = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(6.7), Inches(1.4))
    tf = cap.text_frame; tf.word_wrap = True
    for txt, bold in [("交換cycle-cover(NP困難)でQUBOは貪欲を全インスタンス(8/8)で上回り、厳密解可能域で厳密最適に一致。差は規模で拡大。", True),
                      ("単一臓器・多臓器割当は古典で最適＝量子優位なし（正直な切り分け）→", False)]:
        p = tf.add_paragraph() if tf.paragraphs[0].runs else tf.paragraphs[0]
        r = p.add_run(); r.text = txt; r.font.size = Pt(10.5); r.font.name = JP; r.font.bold = bold
        r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # right panel callouts
    panels = [
        ("割当は分離可能（量子は不要）", "単一臓器=Hungarian最適。多臓器同時割当も救命数では smart greedy=QUBO=厳密最適。当初の複合移植優位仮説はベンチで否定（正直な負の結果）", RGBColor(0x6B, 0x6B, 0x6B)),
        ("FHE暗号計算（実測）", "composite_score 21µs ／ full pipeline 50.91ms/ペア（production_8192, N=8192）。MKFHE bootstrappingで異鍵暗号文の比較を復号なしで実証", RGBColor(0x2E, 0x5B, 0x88)),
        ("スケール", "全臓器スコアリング 国内11分 → GPU 48秒、全世界18分（冷阻血時間内）。N=8192でCUDA 14x高速化", RGBColor(0x4A, 0x90, 0xA4)),
    ]
    y = 1.5
    for title, body, col in panels:
        _box(slide, 7.45, y, 5.35, 1.5, title, body, col, bsz=10.5)
        y += 1.62


def build_software_slide(prs, layout):
    slide = _new_slide(prs, layout, "実装したソフトウェア（実装ステータス）")
    GREEN = RGBColor(0x3A, 0x6B, 0x35)
    AMBER = RGBColor(0xB8, 0x7A, 0x1E)
    GRAY = RGBColor(0x70, 0x70, 0x70)
    rows = [
        ("hyde（PQC / TPM）", "✅ 実装済",
         "hyde / hyde-core / hyde-tpm / hyde-wasm 等6クレート、48テスト。本物の ML-KEM-768＋ML-DSA署名＋TPM封印(aes-gcm)で個人情報を暗号学的に分離", GREEN),
        ("plat（FHE）", "✅ 実装済・実測あり",
         "plat / plat-core / plat-mkfhe / plat-bootstrap / plat-gpu の5クレート、77テスト。bootstrapping・MKFHE・GPU(CUDA 14x)。50.91ms/ペア（production_8192）実測の出所", GREEN),
        ("argo（ZKP）", "✅ 実装済",
         "argo / argo-core、27テスト。Pedersenコミット(Ristretto255)＋Schnorr Σプロトコル(Fiat-Shamir)で知識・一致・線形関係を証明。range proofは今後", GREEN),
        ("niobi（統合・最適化・アプリ）", "✅ 実装済",
         "QUBO・Hungarian・多臓器・ペア交換・5領域example・WASM。58テスト。FHEスコアは実plat-mkfhe、適合性ZKPは実argoを使用。hyde連携とE2E暗号化境界の本番統合が次段階", GREEN),
    ]
    y = 1.55
    for name, status, body, col in rows:
        # status chip
        chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.55), Inches(y), Inches(3.5), Inches(1.12))
        chip.fill.solid(); chip.fill.fore_color.rgb = col
        chip.line.fill.background(); chip.shadow.inherit = False
        tf = chip.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.12)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = name; r.font.size = Pt(12); r.font.bold = True
        r.font.name = JP; r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = status
        r2.font.size = Pt(10.5); r2.font.name = JP; r2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        # body box
        _box(slide, 4.2, y, 8.6, 1.12, "", body, RGBColor(0xF4, 0xF4, 0xF4),
             tcol=RGBColor(0x22, 0x22, 0x22), bsz=10.5)
        y += 1.26
    cap = slide.shapes.add_textbox(Inches(0.55), Inches(6.7), Inches(12.2), Inches(0.5))
    r = cap.text_frame.paragraphs[0].add_run()
    r.text = "暗号3基盤 hyde(PQC)/plat(FHE)/argo(ZKP) は実装・テスト済（計152テスト）。niobi が統合（FHEスコア＝実plat-mkfhe、適合性証明＝実argo）。残るは hyde 連携と E2E 暗号化境界の本番統合。"
    r.font.size = Pt(10.5); r.font.name = JP; r.font.italic = True; r.font.color.rgb = RGBColor(0x44, 0x44, 0x44)


def add_free_slides(prs):
    layout = prs.slides[1].slide_layout  # タイトルとコンテンツ
    build_comparison_slide(prs, layout)
    build_benchmark_slide(prs, layout)
    build_software_slide(prs, layout)


def main():
    prs = Presentation(str(TEMPLATE))
    slides = list(prs.slides)

    # ---- Slide 1: table + delete 作成要領 box ----
    s1 = slides[0]
    order = ["応募代表者氏名", "チーム名", "選択課題ID", "選択課題名", "提案名", "解決案ID"]
    to_delete = []
    for sh in s1.shapes:
        if sh.has_table:
            for ri, key in enumerate(order):
                cell = sh.table.rows[ri].cells[1]
                tf = cell.text_frame
                _clear_tf(tf)
                _add_line(tf, COVER.get(key, ""), 0, 12)
        elif sh.has_text_frame and "作成要領" in sh.text_frame.text:
            to_delete.append(sh)
    for sh in to_delete:
        sh._element.getparent().remove(sh._element)

    # ---- Slides 2-7: fill content placeholders ----
    for si, mapping in FILL.items():
        slide = slides[si]
        by_id = {sh.shape_id: sh for sh in slide.shapes}
        for shape_id, (lines, size) in mapping.items():
            sh = by_id.get(shape_id)
            if sh is None or not sh.has_text_frame:
                print(f"WARN slide{si+1} shape{shape_id} missing")
                continue
            fill_text_frame(sh.text_frame, lines, size)

    # ---- Slide 3: architecture diagram ----
    build_slide3_diagram(slides[2])

    # ---- Free slides 8-10 ----
    add_free_slides(prs)

    # ---- Anonymize document metadata (template carries 作成者 'Hiroki Okuda (JP)') ----
    cp = prs.core_properties
    cp.author = ""
    cp.last_modified_by = ""

    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}  ({len(slides)} slides)")


if __name__ == "__main__":
    main()
